from __future__ import annotations

import hashlib
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import chromadb
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader

try:
    import pandas as pd
except Exception:  # pragma: no cover
    pd = None

from storage.notebook_store import NotebookStore

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
COLLECTION_NAME = "sources"


@dataclass
class IngestResult:
    source_name: str
    num_chunks: int
    status: str
    detail: str


def _embedding_fn():
    # Lazy import to keep Space startup memory low.
    from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction

    return SentenceTransformerEmbeddingFunction(model_name=EMBED_MODEL)


def _chroma_collection(chroma_dir: Path):
    client = chromadb.PersistentClient(path=str(chroma_dir))
    return client.get_or_create_collection(name=COLLECTION_NAME, embedding_function=_embedding_fn())


def _extract_text_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    texts = []
    for page in reader.pages:
        texts.append(page.extract_text() or "")
    return "\n".join(texts)


def _extract_text_from_pptx(path: Path) -> str:
    if Presentation is None:
        return "python-pptx is unavailable in this runtime."
    pres = Presentation(str(path))
    parts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_text_from_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_text_from_xlsx(path: Path) -> str:
    if pd is None:
        raise RuntimeError("pandas/openpyxl not available for .xlsx ingestion")

    workbook = pd.read_excel(path, sheet_name=None)
    parts: list[str] = []
    for sheet_name, frame in workbook.items():
        parts.append(f"## Sheet: {sheet_name}")
        frame = frame.fillna("")
        for row in frame.values.tolist():
            row_text = " | ".join([str(cell).strip() for cell in row if str(cell).strip()])
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_text_from_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_text_from_pdf(path)
    if suffix == ".pptx":
        return _extract_text_from_pptx(path)
    if suffix in {".txt", ".md", ".csv"}:
        return _extract_text_from_txt(path)
    if suffix == ".xlsx":
        return _extract_text_from_xlsx(path)
    raise ValueError(f"Unsupported file type: {suffix}. Use .pdf, .pptx, .txt, .csv, or .xlsx")


def extract_text_from_url(url: str, timeout: int = 20) -> str:
    response = requests.get(url, timeout=timeout)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    cleaned = "\n".join([line.strip() for line in text.splitlines() if line.strip()])
    return cleaned


def chunk_text(text: str, chunk_size: int = 900, overlap: int = 150) -> list[str]:
    stripped = text.strip()
    if not stripped:
        return []

    chunks: list[str] = []
    start = 0
    text_len = len(stripped)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk = stripped[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == text_len:
            break
        start = max(0, end - overlap)
    return chunks


def _write_extracted_text(store: NotebookStore, username: str, notebook_id: str, source_name: str, text: str) -> Path:
    paths = store.notebook_paths(username, notebook_id)
    paths.files_extracted.mkdir(parents=True, exist_ok=True)
    base = Path(source_name).stem or "source"
    extracted_path = paths.files_extracted / f"{base}.txt"
    extracted_path.write_text(text, encoding="utf-8")
    return extracted_path


def _persist_chunks(
    store: NotebookStore,
    username: str,
    notebook_id: str,
    source_name: str,
    source_type: str,
    chunks: list[str],
) -> int:
    paths = store.notebook_paths(username, notebook_id)
    paths.chroma.mkdir(parents=True, exist_ok=True)

    collection = _chroma_collection(paths.chroma)
    source_hash = hashlib.sha1(f"{source_type}:{source_name}".encode("utf-8")).hexdigest()[:12]

    docs = []
    metadatas = []
    ids = []
    for idx, chunk in enumerate(chunks):
        docs.append(chunk)
        ids.append(f"{source_hash}-{idx}")
        metadatas.append(
            {
                "source_name": source_name,
                "source_type": source_type,
                "chunk_index": idx,
            }
        )

    if docs:
        collection.upsert(ids=ids, documents=docs, metadatas=metadatas)
    store.touch_notebook(username, notebook_id)
    return len(docs)


def ingest_file(store: NotebookStore, username: str, notebook_id: str, local_file_path: str) -> IngestResult:
    source_path = Path(local_file_path)
    paths = store.notebook_paths(username, notebook_id)
    paths.files_raw.mkdir(parents=True, exist_ok=True)

    destination = paths.files_raw / source_path.name
    if source_path.resolve() != destination.resolve():
        shutil.copy2(source_path, destination)

    try:
        text = extract_text_from_file(destination)
        chunks = chunk_text(text)
        _write_extracted_text(store, username, notebook_id, destination.name, text)
        count = _persist_chunks(store, username, notebook_id, destination.name, "file", chunks)
        return IngestResult(source_name=destination.name, num_chunks=count, status="ok", detail="Ingested")
    except Exception as exc:
        return IngestResult(source_name=destination.name, num_chunks=0, status="error", detail=str(exc))


def ingest_url(store: NotebookStore, username: str, notebook_id: str, url: str) -> IngestResult:
    parsed = urlparse(url)
    source_name = parsed.netloc + parsed.path
    if not parsed.scheme:
        return IngestResult(source_name=url, num_chunks=0, status="error", detail="Invalid URL")

    try:
        text = extract_text_from_url(url)
        chunks = chunk_text(text)
        _write_extracted_text(store, username, notebook_id, f"url_{hashlib.sha1(url.encode()).hexdigest()[:8]}.txt", text)
        count = _persist_chunks(store, username, notebook_id, source_name, "url", chunks)
        return IngestResult(source_name=source_name, num_chunks=count, status="ok", detail="Ingested")
    except Exception as exc:
        return IngestResult(source_name=source_name, num_chunks=0, status="error", detail=str(exc))


def ingest_many_files(store: NotebookStore, username: str, notebook_id: str, files: list[Any]) -> list[IngestResult]:
    results: list[IngestResult] = []
    for file_obj in files:
        file_path = getattr(file_obj, "name", None) or str(file_obj)
        results.append(ingest_file(store, username, notebook_id, file_path))
    return results
